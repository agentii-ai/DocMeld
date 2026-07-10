"""python-pptx backend — slide-level shape extraction for .pptx files.

Extracts a rich set of element types from PowerPoint presentations:
text, title, table, image, chart, formula, smartart, notes, group,
footer, and comment. Speaker notes are emitted after slide content.
Hidden slides are flagged with ``hidden: true`` and numbered continuously.
"""

from __future__ import annotations

import base64
import logging
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from typing import Any

logger = logging.getLogger("docmeld")

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_DIAGRAM_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"


def _formula_from_shape(shape: Any) -> str | None:
    """Extract OMML math text from a shape, if present.

    Returns the concatenated text of the shape's ``m:oMath`` nodes, or
    ``None`` when the shape contains no embedded equation. This is a
    best-effort textual representation (not full LaTeX conversion).
    """
    try:
        el = shape._element
    except Exception:
        return None
    if el is None:
        return None
    maths = el.findall(f".//{{{_M_NS}}}oMath")
    if not maths:
        return None
    parts: list[str] = []
    for math in maths:
        for t in math.iter(f"{{{_M_NS}}}t"):
            if t.text and t.text.strip():
                parts.append(t.text.strip())
    joined = " ".join(parts).strip()
    return joined or None


def _chart_type_label(chart: Any) -> str:
    """Map a python-pptx chart_type enum to a simple label."""
    try:
        name = chart.chart_type.name.upper()
    except Exception:
        return "unknown"
    if "PIE" in name or "DOUGHNUT" in name:
        return "doughnut" if "DOUGHNUT" in name else "pie"
    if "LINE" in name:
        return "line"
    if "BAR" in name or "COLUMN" in name:
        return "bar"
    if "AREA" in name:
        return "area"
    if "SCATTER" in name or "XY" in name:
        return "scatter"
    if "RADAR" in name:
        return "radar"
    if "BUBBLE" in name:
        return "bubble"
    return "unknown"


def _chart_to_markdown(chart: Any) -> str:
    """Build a markdown table from a chart's categories and series."""
    try:
        plot = chart.plots[0]
        categories = [str(c) for c in plot.categories]
        series = list(plot.series)
    except Exception:
        return ""
    if not series:
        return ""
    headers = ["Category"] + [str(s.name or f"Series {i + 1}") for i, s in enumerate(series)]
    lines = ["| " + " | ".join(headers) + " |"]
    lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
    n = len(categories) if categories else max((len(list(s.values)) for s in series), default=0)
    for r in range(n):
        cat = categories[r] if r < len(categories) else str(r + 1)
        row = [cat]
        for s in series:
            vals = list(s.values)
            row.append(str(vals[r]) if r < len(vals) else "")
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _table_to_markdown(table: Any) -> str:
    """Convert a python-pptx table to a markdown table string.

    Returns an empty string when the table has no rows or every cell is
    empty (e.g. a blank placeholder table), so callers can skip it.
    """
    rows = list(table.rows)
    if not rows:
        return ""
    cell_texts = [[cell.text.replace("\n", " ").strip() for cell in row.cells] for row in rows]
    if not any(any(cells) for cells in cell_texts):
        return ""
    lines: list[str] = []
    for r_idx, cells in enumerate(cell_texts):
        lines.append("| " + " | ".join(cells) + " |")
        if r_idx == 0:
            lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(lines)


def _text_with_hyperlinks(text_frame: Any) -> str:
    """Extract text from a text frame, preserving hyperlinks as markdown."""
    paragraphs: list[str] = []
    for para in text_frame.paragraphs:
        run_parts: list[str] = []
        for run in para.runs:
            txt = run.text or ""
            addr = None
            try:
                addr = run.hyperlink.address
            except Exception:
                addr = None
            if addr:
                run_parts.append(f"[{txt}]({addr})")
            else:
                run_parts.append(txt)
        line = "".join(run_parts)
        if line.strip():
            paragraphs.append(line)
    return "\n".join(paragraphs)


def _extract_comments_by_slide(pptx_path: str) -> dict[int, list[tuple[str, str]]]:
    """Parse legacy PPTX comment parts, keyed by 1-based slide number.

    Returns a mapping slide_no -> list of (author, text). Best-effort:
    missing or malformed comment parts are skipped without error.
    """
    result: dict[int, list[tuple[str, str]]] = {}
    try:
        with zipfile.ZipFile(pptx_path) as z:
            names = set(z.namelist())
            # Author id -> name
            authors: dict[str, str] = {}
            if "ppt/commentAuthors.xml" in names:
                try:
                    root = ET.fromstring(z.read("ppt/commentAuthors.xml"))
                    for a in root.findall(f"{{{_P_NS}}}cmAuthor"):
                        authors[a.get("id", "")] = a.get("name", "")
                except ET.ParseError:
                    pass
            # Each slide's rels may point to a comments part
            slide_names = [
                n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            ]
            for slide_name in slide_names:
                stem = Path(slide_name).stem  # slideN
                try:
                    slide_no = int(stem.replace("slide", ""))
                except ValueError:
                    continue
                rels_name = f"ppt/slides/_rels/{stem}.xml.rels"
                if rels_name not in names:
                    continue
                try:
                    rels = ET.fromstring(z.read(rels_name))
                except ET.ParseError:
                    continue
                for rel in rels.findall(f"{{{_R_NS}}}Relationship"):
                    if not (rel.get("Type") or "").endswith("/comments"):
                        continue
                    target = rel.get("Target", "").replace("../", "ppt/")
                    if target not in names:
                        continue
                    try:
                        cm_root = ET.fromstring(z.read(target))
                    except ET.ParseError:
                        continue
                    for cm in cm_root.findall(f"{{{_P_NS}}}cm"):
                        author = authors.get(cm.get("authorId", ""), "")
                        text_el = cm.find(f"{{{_P_NS}}}text")
                        text = (text_el.text or "") if text_el is not None else ""
                        if text.strip():
                            result.setdefault(slide_no, []).append((author, text.strip()))
    except (zipfile.BadZipFile, KeyError, OSError):
        pass
    return result


def _smartart_text(shape: Any, pptx_path: str) -> str | None:
    """Best-effort SmartArt text extraction from the diagram data part."""
    try:
        el = shape._element
        # Find graphicData with the diagram URI
        gd = el.find(f".//{{{_A_NS}}}graphicData")
        if gd is None or _DIAGRAM_URI not in (gd.get("uri") or ""):
            return None
    except Exception:
        return None
    # Read all <a:t> text nodes from any diagram data part in the package
    try:
        with zipfile.ZipFile(pptx_path) as z:
            data_parts = [n for n in z.namelist() if "diagrams/data" in n and n.endswith(".xml")]
            texts: list[str] = []
            for part in data_parts:
                try:
                    root = ET.fromstring(z.read(part))
                except ET.ParseError:
                    continue
                for t in root.iter(f"{{{_A_NS}}}t"):
                    if t.text and t.text.strip():
                        texts.append(t.text.strip())
            if texts:
                return "\n".join(f"- {t}" for t in texts)
    except (zipfile.BadZipFile, OSError):
        return None
    return None


class PptxBackend:
    """Extract elements from .pptx presentations using python-pptx."""

    def extract_elements(
        self, pptx_path: str, output_dir: str  # noqa: ARG002
    ) -> list[dict[str, Any]]:
        """Extract slide elements from a .pptx file.

        Args:
            pptx_path: Path to the .pptx file.
            output_dir: Directory for auxiliary outputs (unused; images are inlined).

        Returns:
            Ordered list of element dicts (without element_id/parent_id, which are
            assigned by shared post-processing).
        """
        from pptx import Presentation

        self._gid_counter = 0
        path_obj = Path(pptx_path)
        if path_obj.suffix.lower() != ".pptx":
            msg = f"PptxBackend only supports .pptx files, got: {path_obj.suffix}"
            raise ValueError(msg)

        try:
            prs = Presentation(pptx_path)
        except Exception as exc:
            msg = f"Failed to open .pptx (corrupt or password-protected?): {path_obj.name}: {exc}"
            raise RuntimeError(msg) from exc

        comments_by_slide = _extract_comments_by_slide(pptx_path)
        all_elements: list[dict[str, Any]] = []
        img_counter = 0

        for idx, slide in enumerate(prs.slides):
            page_no = idx + 1
            hidden = False
            try:
                hidden = slide._element.get("show") == "0"
            except Exception:
                hidden = False

            sortable: list[tuple[int, int, int, dict[str, Any]]] = []
            for z, shape in enumerate(slide.shapes, start=1):
                top = int(getattr(shape, "top", None) or 0)
                left = int(getattr(shape, "left", None) or 0)
                for elem in self._map_shape(shape, page_no, hidden, pptx_path):
                    sortable.append((top, left, z, elem))
                    img_counter = self._maybe_number_image(elem, img_counter)

            # Hybrid ordering: geometric (top, left) primary, z-order tie-break
            sortable.sort(key=lambda t: (t[0], t[1], t[2]))
            slide_elements = [t[3] for t in sortable]

            if len(slide_elements) > 20:
                logger.debug(f"Slide {page_no}: complex layering ({len(slide_elements)} elements)")

            # Comments (anchored to slide)
            for author, text in comments_by_slide.get(page_no, []):
                slide_elements.append(
                    {
                        "type": "comment",
                        "content": text,
                        "author": author,
                        "page_no": page_no,
                        "hidden": hidden,
                    }
                )

            # Speaker notes — emitted after all slide content
            if slide.has_notes_slide:
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                except Exception:
                    notes_text = ""
                if notes_text:
                    slide_elements.append(
                        {
                            "type": "notes",
                            "content": notes_text,
                            "page_no": page_no,
                            "hidden": hidden,
                        }
                    )

            all_elements.extend(slide_elements)

        logger.info(
            f"Pptx: {path_obj.name} → {len(all_elements)} elements, {len(prs.slides._sldIdLst)} slides"
        )
        return all_elements

    def _maybe_number_image(self, elem: dict[str, Any], counter: int) -> int:
        if elem.get("type") == "image" and not elem.get("image_id"):
            counter += 1
            elem["image_id"] = f"pptx_image_{counter:04d}"
            if not elem.get("image_name"):
                elem["image_name"] = f"{elem['image_id']}.png"
        return counter

    def _map_shape(
        self, shape: Any, page_no: int, hidden: bool, pptx_path: str, depth: int = 0
    ) -> list[dict[str, Any]]:
        """Map a single pptx shape to zero or more element dicts."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

        elements: list[dict[str, Any]] = []

        if depth > 5:
            logger.warning(
                f"Slide {page_no}: nested shape depth exceeds 5; flattening deeper content"
            )

        # Exclude non-semantic media / embedded OLE objects with a warning
        try:
            stype = shape.shape_type
        except Exception:
            stype = None
        if stype in (
            getattr(MSO_SHAPE_TYPE, "MEDIA", None),
            getattr(MSO_SHAPE_TYPE, "EMBEDDED_OLE_OBJECT", None),
            getattr(MSO_SHAPE_TYPE, "OLE_CONTROL_OBJECT", None),
        ):
            logger.warning(
                f"Slide {page_no}: skipping media/OLE object "
                f"'{getattr(shape, 'name', '')}' (not extractable)"
            )
            return elements

        # Group shape → group element + recurse children
        try:
            is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if is_group:
            children = list(shape.shapes)
            self._gid_counter += 1
            gid = self._gid_counter
            elements.append(
                {
                    "type": "group",
                    "content": f"Group of {len(children)} shapes",
                    "child_count": len(children),
                    "page_no": page_no,
                    "hidden": hidden,
                    "_gid": gid,
                }
            )
            for child in children:
                child_elems = self._map_shape(child, page_no, hidden, pptx_path, depth + 1)
                for ce in child_elems:
                    ce.setdefault("_pgid", gid)
                elements.extend(child_elems)
            return elements

        # Table
        if getattr(shape, "has_table", False):
            md = _table_to_markdown(shape.table)
            if md:
                elements.append(
                    {
                        "type": "table",
                        "content": md,
                        "summary": "",
                        "page_no": page_no,
                        "hidden": hidden,
                    }
                )
            return elements

        # Chart
        if getattr(shape, "has_chart", False):
            chart = shape.chart
            md = _chart_to_markdown(chart)
            ctype = _chart_type_label(chart)
            if md:
                elements.append(
                    {
                        "type": "chart",
                        "chart_type": ctype,
                        "content": md,
                        "image": "",
                        "image_name": f"chart_{page_no}.png",
                        "page_no": page_no,
                        "hidden": hidden,
                    }
                )
            else:
                elements.append(
                    {
                        "type": "image",
                        "content": f"chart ({ctype}); data unavailable",
                        "image": "",
                        "image_id": "",
                        "image_name": "",
                        "bbox": (0.0, 0.0, 0.0, 0.0),
                        "page_no": page_no,
                        "hidden": hidden,
                    }
                )
            return elements

        # SmartArt (diagram graphic frame)
        smart = _smartart_text(shape, pptx_path)
        if smart is not None:
            elements.append(
                {
                    "type": "smartart",
                    "smartart_type": "unknown",
                    "content": smart,
                    "image": "",
                    "image_name": f"smartart_{page_no}.png",
                    "page_no": page_no,
                    "hidden": hidden,
                }
            )
            return elements

        # Picture
        try:
            is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        except Exception:
            is_picture = False
        if is_picture:
            try:
                image = shape.image
                b64 = base64.b64encode(image.blob).decode("utf-8")
                ext = image.ext or "png"
                bbox = (
                    float(getattr(shape, "left", None) or 0),
                    float(getattr(shape, "top", None) or 0),
                    float(getattr(shape, "width", None) or 0),
                    float(getattr(shape, "height", None) or 0),
                )
                elements.append(
                    {
                        "type": "image",
                        "content": "",
                        "image": b64,
                        "image_id": "",
                        "image_name": f"pptx_image.{ext}",
                        "bbox": bbox,
                        "page_no": page_no,
                        "hidden": hidden,
                    }
                )
            except Exception:
                pass
            return elements

        # Text-bearing shapes (placeholders, text boxes)
        if getattr(shape, "has_text_frame", False):
            # Embedded OMML equation → formula element (best-effort text)
            formula = _formula_from_shape(shape)
            if formula:
                elements.append(
                    {
                        "type": "formula",
                        "content": formula,
                        "formula_type": "OMML",
                        "page_no": page_no,
                        "hidden": hidden,
                    }
                )

            text = _text_with_hyperlinks(shape.text_frame)
            if not text.strip():
                return elements

            ph_type = None
            if getattr(shape, "is_placeholder", False):
                try:
                    ph_type = shape.placeholder_format.type
                except Exception:
                    ph_type = None

            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                elements.append(
                    {
                        "type": "title",
                        "level": 0,
                        "content": text.strip(),
                        "page_no": page_no,
                        "hidden": hidden,
                    }
                )
            elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                elements.append(
                    {
                        "type": "title",
                        "level": 1,
                        "content": text.strip(),
                        "page_no": page_no,
                        "hidden": hidden,
                    }
                )
            elif ph_type in (
                PP_PLACEHOLDER.FOOTER,
                PP_PLACEHOLDER.SLIDE_NUMBER,
                PP_PLACEHOLDER.DATE,
            ):
                elements.append(
                    {
                        "type": "footer",
                        "content": text.strip(),
                        "page_scope": "all",
                        "page_no": page_no,
                        "hidden": hidden,
                    }
                )
            else:
                elements.append(
                    {"type": "text", "content": text, "page_no": page_no, "hidden": hidden}
                )
        return elements

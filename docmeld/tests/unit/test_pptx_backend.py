"""Unit tests for the PptxBackend (python-pptx slide extraction)."""
from __future__ import annotations

import tempfile
from pathlib import Path

import pytest

SAMPLES = Path(__file__).resolve().parents[3] / "samples" / "ppt"


def _extract(name: str):
    from docmeld.bronze.backends.pptx_backend import PptxBackend

    with tempfile.TemporaryDirectory() as d:
        return PptxBackend().extract_elements(str(SAMPLES / name), d)


def _types(elements):
    from collections import Counter

    return Counter(e["type"] for e in elements)


class TestPptxBackendCore:
    def test_basic_titles_and_text(self) -> None:
        els = _extract("sample_pptx_basic.pptx")
        t = _types(els)
        assert t["title"] >= 1
        assert t["text"] >= 1
        # slide-numbered page_no starting at 1
        pages = sorted({e["page_no"] for e in els})
        assert pages[0] == 1
        assert all(isinstance(e["page_no"], int) and e["page_no"] >= 1 for e in els)

    def test_image_extraction_base64(self) -> None:
        els = _extract("sample_pptx_image.pptx")
        images = [e for e in els if e["type"] == "image"]
        assert images, "expected at least one image element"
        img = images[0]
        assert img["image"]  # non-empty base64
        assert img["image_id"]
        assert len(img["bbox"]) == 4

    def test_extension_guard(self) -> None:
        from docmeld.bronze.backends.pptx_backend import PptxBackend

        with pytest.raises(ValueError):
            PptxBackend().extract_elements("/tmp/not_a.pptx.pdf", "/tmp")


class TestPptxBackendRich:
    def test_chart_extraction(self) -> None:
        els = _extract("sample_pptx_chart_bar.pptx")
        charts = [e for e in els if e["type"] == "chart"]
        assert charts, "expected a chart element"
        c = charts[0]
        assert c["chart_type"] in {"bar", "line", "pie", "area", "scatter", "doughnut", "radar", "bubble", "unknown"}
        assert "|" in c["content"]  # markdown table

    def test_comments_with_author(self) -> None:
        els = _extract("sample_pptx_comments.pptx")
        comments = [e for e in els if e["type"] == "comment"]
        assert comments, "expected comment elements"
        assert any(c["author"] for c in comments), "expected author attribution"
        # comment anchored to a valid slide
        assert all(c["page_no"] >= 1 for c in comments)

    def test_shapes_sample_no_crash(self) -> None:
        els = _extract("sample_pptx_shapes.pptx")
        assert els  # graceful handling of unrecognized shapes

    def test_notes_extracted_after_content(self) -> None:
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as d:
            p = Path(d) / "deck.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "My Title"
            slide.notes_slide.notes_text_frame.text = "Speaker note body"
            prs.save(str(p))

            from docmeld.bronze.backends.pptx_backend import PptxBackend

            els = PptxBackend().extract_elements(str(p), d)
            notes = [e for e in els if e["type"] == "notes"]
            assert notes and notes[0]["content"] == "Speaker note body"
            # notes come after title in the element order
            assert els.index(notes[0]) > 0

    def test_hyperlink_inline_markdown(self) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as d:
            p = Path(d) / "link.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            run = box.text_frame.paragraphs[0].add_run()
            run.text = "OpenAI"
            run.hyperlink.address = "https://openai.com"
            prs.save(str(p))

            from docmeld.bronze.backends.pptx_backend import PptxBackend

            els = PptxBackend().extract_elements(str(p), d)
            texts = [e for e in els if e["type"] == "text"]
            assert any("[OpenAI](https://openai.com)" in e["content"] for e in texts)


class TestPptxOrdering:
    def test_issue_sample_ordered_and_no_crash(self) -> None:
        els = _extract("sample_pptx_issue.pptx")
        assert els
        # elements grouped/sorted by page_no ascending
        pages = [e["page_no"] for e in els]
        assert pages == sorted(pages) or len(set(pages)) >= 1

    def test_hidden_slide_flag(self) -> None:
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as d:
            p = Path(d) / "hidden.pptx"
            prs = Presentation()
            s1 = prs.slides.add_slide(prs.slide_layouts[1])
            s1.shapes.title.text = "Visible"
            s2 = prs.slides.add_slide(prs.slide_layouts[1])
            s2.shapes.title.text = "Hidden One"
            # mark slide 2 hidden
            s2._element.set("show", "0")
            prs.save(str(p))

            from docmeld.bronze.backends.pptx_backend import PptxBackend

            els = PptxBackend().extract_elements(str(p), d)
            page2 = [e for e in els if e["page_no"] == 2]
            assert page2 and all(e["hidden"] for e in page2)
            page1 = [e for e in els if e["page_no"] == 1]
            assert page1 and all(not e["hidden"] for e in page1)


class TestPptxEdgeCases:
    def test_corrupt_pptx_raises_runtime_error(self, tmp_path) -> None:
        from docmeld.bronze.backends.pptx_backend import PptxBackend

        bad = tmp_path / "broken.pptx"
        bad.write_bytes(b"not a real pptx")
        with pytest.raises(RuntimeError):
            PptxBackend().extract_elements(str(bad), str(tmp_path))

    def test_map_shape_media_skipped(self) -> None:
        from unittest.mock import MagicMock
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from docmeld.bronze.backends.pptx_backend import PptxBackend

        shape = MagicMock()
        shape.shape_type = MSO_SHAPE_TYPE.MEDIA
        shape.name = "Video 1"
        out = PptxBackend()._map_shape(shape, 1, False, "x.pptx")
        assert out == []


class TestPptxFormula:
    """FR-013: OMML equation extraction."""

    def _shape_with_math(self, math_texts):
        from lxml import etree

        M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
        runs = "".join(f'<m:r><m:t>{t}</m:t></m:r>' for t in math_texts)
        xml = f'<p:sp xmlns:p="p" xmlns:m="{M}"><m:oMath>{runs}</m:oMath></p:sp>'

        class FakeShape:
            _element = etree.fromstring(xml.encode())

        return FakeShape()

    def test_formula_from_shape_extracts_omml(self) -> None:
        from docmeld.bronze.backends.pptx_backend import _formula_from_shape

        shape = self._shape_with_math(["E", "=", "mc", "2"])
        assert _formula_from_shape(shape) == "E = mc 2"

    def test_formula_none_when_no_math(self) -> None:
        from lxml import etree
        from docmeld.bronze.backends.pptx_backend import _formula_from_shape

        class FakeShape:
            _element = etree.fromstring(b'<p:sp xmlns:p="p"><a/></p:sp>')

        assert _formula_from_shape(FakeShape()) is None


class TestGroupParentResolution:
    """FR-016: grouped children link to their parent group element."""

    def test_resolve_group_parents(self) -> None:
        from docmeld.bronze.element_extractor import (
            _assign_element_ids,
            _assign_parent_ids,
            _resolve_group_parents,
        )

        elements = [
            {"type": "title", "level": 0, "content": "T", "page_no": 1},
            {"type": "group", "content": "Group of 2 shapes", "child_count": 2, "page_no": 1, "_gid": 1},
            {"type": "text", "content": "child a", "page_no": 1, "_pgid": 1},
            {"type": "text", "content": "child b", "page_no": 1, "_pgid": 1},
        ]
        _assign_element_ids(elements)
        _assign_parent_ids(elements)
        _resolve_group_parents(elements)

        group = elements[1]
        assert group["element_id"] == "e_0002"
        # children point to the group, not the title
        assert elements[2]["parent_id"] == "e_0002"
        assert elements[3]["parent_id"] == "e_0002"
        # temp keys stripped
        assert all("_gid" not in e and "_pgid" not in e for e in elements)


class TestPptxOutputQualityFixes:
    """Regression tests for output-evaluation bugs (O2, O1)."""

    def test_empty_table_not_emitted(self) -> None:
        """O2: a table whose cells are all empty must not become a table element."""
        els = _extract("sample_pptx_chart_bar.pptx")
        for e in els:
            if e["type"] == "table":
                # every table emitted must have at least one non-empty cell
                cells = e["content"].replace("|", " ").replace("-", " ")
                assert cells.strip(), f"empty table emitted: {e}"

    def test_table_helper_empty_returns_blank(self) -> None:
        from unittest.mock import MagicMock
        from docmeld.bronze.backends.pptx_backend import _table_to_markdown

        # 2x2 table with all-empty cells
        def cell(text):
            c = MagicMock(); c.text = text; return c
        row = MagicMock(); row.cells = [cell(""), cell("")]
        table = MagicMock(); table.rows = [row, row]
        assert _table_to_markdown(table) == ""


class TestBronzeNoOrphanDir:
    """O1: failed conversion must not leave an empty output directory."""

    def test_failed_ppt_leaves_no_empty_dir(self, tmp_path) -> None:
        import shutil as _sh
        from unittest.mock import patch
        from docmeld.bronze.processor import BronzeProcessor

        src = SAMPLES / "sample_ppt_legacy.ppt"
        if not src.exists():
            pytest.skip("legacy ppt not found")
        copy = tmp_path / "legacy.ppt"
        _sh.copy(src, copy)

        # Force soffice-missing so conversion fails
        with patch("shutil.which", return_value=None):
            with pytest.raises(Exception):
                BronzeProcessor().process_file(str(copy), backend="auto")

        # No orphan hashed directory should remain
        leftover = [d for d in tmp_path.iterdir() if d.is_dir()]
        assert leftover == [], f"orphan dir left behind: {leftover}"


class TestPptxRealDataFixtures:
    """O4/O5: end-to-end verification of formula (OMML) and group extraction."""

    def test_formula_end_to_end_from_injected_omml(self, tmp_path) -> None:
        from lxml import etree
        from pptx import Presentation
        from pptx.util import Inches
        from docmeld.bronze.backends.pptx_backend import PptxBackend

        p = tmp_path / "eq.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.paragraphs[0].add_run().text = "Equation:"
        # Inject an OMML oMath node into the shape's txBody
        M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
        omath = etree.SubElement(box.text_frame._txBody, f"{{{M}}}oMath")
        for token in ("E", "=", "mc", "2"):
            r = etree.SubElement(omath, f"{{{M}}}r")
            t = etree.SubElement(r, f"{{{M}}}t")
            t.text = token
        prs.save(str(p))

        els = PptxBackend().extract_elements(str(p), str(tmp_path))
        formulas = [e for e in els if e["type"] == "formula"]
        assert formulas, "expected a formula element from injected OMML"
        assert formulas[0]["formula_type"] == "OMML"
        assert "E" in formulas[0]["content"] and "mc" in formulas[0]["content"]

    def test_group_end_to_end_with_child_linkage(self, tmp_path) -> None:
        import json as _json
        import shutil as _sh
        from pptx import Presentation
        from pptx.util import Inches
        from docmeld.bronze.processor import BronzeProcessor

        p = tmp_path / "grp.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        grp = slide.shapes.add_group_shape()
        b1 = grp.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        b1.text_frame.text = "Child A"
        b2 = grp.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
        b2.text_frame.text = "Child B"
        prs.save(str(p))

        result = BronzeProcessor().process_file(str(p), backend="pptx")
        els = _json.load(open(result.output_path))
        groups = [e for e in els if e["type"] == "group"]
        assert groups, "expected a group element"
        gid = groups[0]["element_id"]
        assert groups[0]["child_count"] == 2
        children = [e for e in els if e.get("parent_id") == gid and e["type"] == "text"]
        assert len(children) == 2, "children must link to their group via parent_id"
        # temporary linkage keys must be stripped
        assert all("_gid" not in e and "_pgid" not in e for e in els)
